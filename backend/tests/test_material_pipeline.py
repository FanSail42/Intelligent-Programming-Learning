import tempfile
from pathlib import Path
from unittest.mock import AsyncMock, patch

import pytest

from app.models.material import CourseMaterial, MaterialStatus, MaterialType
from app.services.chunking import chunk_blocks
from app.services.document_parser import ParsedBlock
from app.services.material_pipeline import process_material
from tests.conftest import login


def test_chunk_blocks_count():
    blocks = [ParsedBlock(text="word " * 600, page=1)]
    pieces = chunk_blocks(blocks, chunk_size=512, overlap=50)
    assert len(pieces) >= 1


@pytest.mark.asyncio
async def test_upload_triggers_pipeline(client, seed_users, seed_course, seed_warehouses, tmp_path):
    content = b"# Python intro\nprint('hello')\n" * 20
    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")

    with patch("app.api.v1.materials._dispatch_process") as mock_dispatch:
        resp = await client.post(
            "/api/v1/materials/upload",
            headers={"Authorization": f"Bearer {tokens['access_token']}"},
            data={"course_id": str(seed_course.id)},
            files={"file": ("demo.md", content, "text/markdown")},
        )
        assert resp.json()["code"] == 0
        mock_dispatch.assert_called_once()


def test_process_material_end_to_end(db_session, seed_course, tmp_path):
    file_path = tmp_path / "demo.md"
    file_path.write_text("# Title\n" + "content line\n" * 50, encoding="utf-8")

    material = CourseMaterial(
        course_id=seed_course.id,
        type=MaterialType.md,
        file_path=str(file_path),
        original_name="demo.md",
        status=MaterialStatus.uploaded,
    )
    db_session.add(material)
    db_session.commit()
    db_session.refresh(material)

    chroma_dir = tempfile.mkdtemp()
    with patch("app.services.material_pipeline.get_vector_store") as mock_vs:
        from app.services.vector_store import VectorStore

        store = VectorStore(persist_dir=chroma_dir)
        mock_vs.return_value = store
        process_material(material.id)

    db_session.refresh(material)
    assert material.status == MaterialStatus.ready


@pytest.mark.asyncio
async def test_retry_failed_material(client, db_session, seed_users, seed_course, tmp_path):
    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")
    headers = {"Authorization": f"Bearer {tokens['access_token']}"}

    file_path = tmp_path / "fail.md"
    file_path.write_text("# fail\n", encoding="utf-8")
    material = CourseMaterial(
        course_id=seed_course.id,
        type=MaterialType.md,
        file_path=str(file_path),
        original_name="fail.md",
        status=MaterialStatus.failed,
        error_message="parse error",
    )
    db_session.add(material)
    db_session.commit()
    db_session.refresh(material)

    with patch("app.api.v1.materials._dispatch_process") as mock_dispatch:
        resp = await client.post(
            f"/api/v1/materials/{material.id}/retry",
            headers=headers,
        )
        assert resp.json()["code"] == 0
        mock_dispatch.assert_called_once()

    db_session.refresh(material)
    assert material.status == MaterialStatus.uploaded
    assert material.error_message is None


@pytest.mark.asyncio
async def test_retry_non_failed_rejected(client, db_session, seed_users, seed_course, tmp_path):
    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")
    headers = {"Authorization": f"Bearer {tokens['access_token']}"}

    file_path = tmp_path / "ok.md"
    file_path.write_text("# ok\n", encoding="utf-8")
    material = CourseMaterial(
        course_id=seed_course.id,
        type=MaterialType.md,
        file_path=str(file_path),
        original_name="ok.md",
        status=MaterialStatus.ready,
    )
    db_session.add(material)
    db_session.commit()
    db_session.refresh(material)

    resp = await client.post(
        f"/api/v1/materials/{material.id}/retry",
        headers=headers,
    )
    assert resp.json()["code"] == 40001


@pytest.mark.asyncio
async def test_delete_material(client, db_session, seed_users, seed_course, tmp_path):
    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")
    headers = {"Authorization": f"Bearer {tokens['access_token']}"}

    file_path = tmp_path / "del.md"
    file_path.write_text("# del\n", encoding="utf-8")
    material = CourseMaterial(
        course_id=seed_course.id,
        type=MaterialType.md,
        file_path=str(file_path),
        original_name="del.md",
        status=MaterialStatus.ready,
    )
    db_session.add(material)
    db_session.commit()
    db_session.refresh(material)

    with patch("app.api.v1.materials.get_vector_store") as mock_vs:
        mock_vs.return_value.delete_by_material.return_value = None
        resp = await client.delete(
            f"/api/v1/materials/{material.id}",
            headers=headers,
        )
        assert resp.json()["code"] == 0
        mock_vs.return_value.delete_by_material.assert_called_once_with(material.id)

    db_session.refresh(material)
    assert material.deleted == 1


@pytest.mark.asyncio
async def test_download_material(client, db_session, seed_users, seed_course, tmp_path):
    from app.api.v1 import materials as materials_api

    materials_api.settings.upload_dir = str(tmp_path)
    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")
    headers = {"Authorization": f"Bearer {tokens['access_token']}"}

    file_path = tmp_path / "slide.pdf"
    file_path.write_bytes(b"%PDF-1.4 download test")
    material = CourseMaterial(
        course_id=seed_course.id,
        type=MaterialType.pdf,
        file_path=str(file_path),
        original_name="slide.pdf",
        status=MaterialStatus.ready,
    )
    db_session.add(material)
    db_session.commit()
    db_session.refresh(material)

    resp = await client.get(
        f"/api/v1/materials/{material.id}/download",
        headers=headers,
    )
    assert resp.status_code == 200
    assert resp.content.startswith(b"%PDF")
    assert "slide.pdf" in resp.headers.get("content-disposition", "")


@pytest.mark.asyncio
async def test_student_cannot_download_material(
    client, db_session, seed_users, seed_course, tmp_path
):
    from app.models.course import CourseStudent

    student = seed_users["student"]
    db_session.add(CourseStudent(course_id=seed_course.id, user_id=student.id))
    file_path = tmp_path / "note.md"
    file_path.write_text("# private", encoding="utf-8")
    material = CourseMaterial(
        course_id=seed_course.id,
        type=MaterialType.md,
        file_path=str(file_path),
        original_name="note.md",
        status=MaterialStatus.ready,
    )
    db_session.add(material)
    db_session.commit()
    db_session.refresh(material)

    tokens = await login(client, student.username, "Student123!")
    resp = await client.get(
        f"/api/v1/materials/{material.id}/download",
        headers={"Authorization": f"Bearer {tokens['access_token']}"},
    )
    assert resp.json()["code"] == 40301


@pytest.mark.asyncio
async def test_upload_duplicate_same_course_rejected(
    client, seed_users, seed_course, seed_warehouses, tmp_path
):
    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")
    headers = {"Authorization": f"Bearer {tokens['access_token']}"}
    content = b"# dup\n" * 10

    with patch("app.api.v1.materials._dispatch_process"):
        resp1 = await client.post(
            "/api/v1/materials/upload",
            headers=headers,
            data={"course_id": str(seed_course.id)},
            files={"file": ("dup.md", content, "text/markdown")},
        )
        assert resp1.json()["code"] == 0

        resp2 = await client.post(
            "/api/v1/materials/upload",
            headers=headers,
            data={"course_id": str(seed_course.id)},
            files={"file": ("dup.md", content, "text/markdown")},
        )
        body = resp2.json()
        assert body["code"] == 40002
        assert "已存在" in body["message"]


@pytest.mark.asyncio
async def test_upload_links_existing_material_other_course(
    client, db_session, seed_users, seed_course, seed_warehouses, tmp_path
):
    from app.models.course import Course

    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")
    headers = {"Authorization": f"Bearer {tokens['access_token']}"}

    other = Course(
        name="Other Course",
        description="",
        teacher_id=teacher.id,
        status="published",
        create_approval="approved",
        publish_approval="approved",
    )
    db_session.add(other)
    db_session.commit()
    db_session.refresh(other)

    content = b"# shared\n" * 10
    with patch("app.api.v1.materials._dispatch_process"):
        resp1 = await client.post(
            "/api/v1/materials/upload",
            headers=headers,
            data={"course_id": str(seed_course.id)},
            files={"file": ("shared.md", content, "text/markdown")},
        )
        assert resp1.json()["code"] == 0

        resp2 = await client.post(
            "/api/v1/materials/upload",
            headers=headers,
            data={"course_id": str(other.id)},
            files={"file": ("shared.md", content, "text/markdown")},
        )
        body = resp2.json()
        assert body["code"] == 0
        assert body["data"]["linked"] is True


@pytest.mark.asyncio
async def test_upload_rejects_oversized_file(client, seed_users, seed_course, seed_warehouses):
    from app.api.v1 import materials as materials_api

    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")
    headers = {"Authorization": f"Bearer {tokens['access_token']}"}

    original_limit = materials_api.settings.max_upload_size_mb
    materials_api.settings.max_upload_size_mb = 1
    try:
        content = b"x" * (1024 * 1024 + 1)
        resp = await client.post(
            "/api/v1/materials/upload",
            headers=headers,
            data={"course_id": str(seed_course.id)},
            files={"file": ("big.txt", content, "text/plain")},
        )
    finally:
        materials_api.settings.max_upload_size_mb = original_limit

    body = resp.json()
    assert body["code"] == 40001
    assert "1MB" in body["message"]


@pytest.mark.asyncio
async def test_upload_rejects_invalid_pdf(client, seed_users, seed_course, seed_warehouses):
    teacher = seed_users["teacher"]
    tokens = await login(client, teacher.username, "Teacher123!")
    headers = {"Authorization": f"Bearer {tokens['access_token']}"}

    resp = await client.post(
        "/api/v1/materials/upload",
        headers=headers,
        data={"course_id": str(seed_course.id)},
        files={"file": ("fake.pdf", b"not a pdf", "application/pdf")},
    )
    body = resp.json()
    assert body["code"] == 40001
    assert "PDF" in body["message"]


def test_process_material_clones_ready_source(db_session, seed_course, tmp_path):
    file_path = tmp_path / "shared.md"
    file_path.write_text("# Shared\n" + "line\n" * 30, encoding="utf-8")

    source = CourseMaterial(
        course_id=seed_course.id,
        type=MaterialType.md,
        file_path=str(file_path),
        original_name="shared.md",
        status=MaterialStatus.ready,
    )
    db_session.add(source)
    db_session.commit()
    db_session.refresh(source)

    chroma_dir = tempfile.mkdtemp()
    with patch("app.services.material_pipeline.get_vector_store") as mock_vs:
        from app.services.vector_store import VectorStore

        store = VectorStore(persist_dir=chroma_dir)
        mock_vs.return_value = store
        process_material(source.id)

    db_session.refresh(source)
    assert source.status == MaterialStatus.ready

    target = CourseMaterial(
        course_id=seed_course.id,
        type=MaterialType.md,
        file_path=str(file_path),
        original_name="shared.md",
        status=MaterialStatus.uploaded,
    )
    db_session.add(target)
    db_session.commit()
    db_session.refresh(target)

    with patch("app.services.material_pipeline.get_vector_store") as mock_vs:
        from app.services.vector_store import VectorStore

        store = VectorStore(persist_dir=chroma_dir)
        mock_vs.return_value = store
        with patch("app.services.material_pipeline.asyncio.run") as mock_run:
            process_material(target.id)
            mock_run.assert_not_called()

    db_session.refresh(target)
    assert target.status == MaterialStatus.ready


def test_parse_pptx_extracts_text(tmp_path):
    pytest.importorskip("pptx")
    from pptx import Presentation

    from app.models.material import MaterialType
    from app.services.document_parser import parse_document

    path = tmp_path / "demo.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Hello PPT"
    prs.save(str(path))

    blocks = parse_document(str(path), MaterialType.pptx)
    assert len(blocks) == 1
    assert "Hello PPT" in blocks[0].text
