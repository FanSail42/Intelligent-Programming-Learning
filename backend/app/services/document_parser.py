from dataclasses import dataclass
from pathlib import Path

import pdfplumber

from app.models.material import MaterialType

PDF_PARSE_EMPTY_MSG = (
    "PDF 未提取到可检索文字。"
    "若由 PPT 转换而来，可能是图片型 PDF（无文字层）。"
    "请使用 PowerPoint「文件→导出→PDF」并勾选可访问性选项，或直接上传 .pptx 文件。"
)

PPTX_PARSE_EMPTY_MSG = (
    "PPTX 中未提取到可检索文字。"
    "若幻灯片主要为图片，系统无法识别图片内文字（不做 OCR）。"
    "请补充可编辑文本、演讲者备注，或使用含文字层的 PDF。"
)


@dataclass
class ParsedBlock:
    text: str
    page: int | None = None


def parse_document(file_path: str, material_type: MaterialType) -> list[ParsedBlock]:
    path = Path(file_path)
    if material_type == MaterialType.pdf:
        return _parse_pdf(path)
    if material_type == MaterialType.pptx:
        return _parse_pptx(path)
    if material_type in (MaterialType.txt, MaterialType.md):
        content = path.read_text(encoding="utf-8", errors="ignore")
        if not content.strip():
            raise ValueError("文件内容为空")
        return [ParsedBlock(text=content, page=None)]
    raise ValueError(f"Unsupported material type: {material_type}")


def _is_meaningful_text(text: str) -> bool:
    stripped = text.strip()
    if len(stripped) < 2:
        return False
    # Ignore pages that are mostly punctuation / whitespace noise.
    alnum = sum(1 for ch in stripped if ch.isalnum())
    return alnum >= max(2, len(stripped) // 20)


def _parse_pdf(path: Path) -> list[ParsedBlock]:
    if not path.is_file():
        raise ValueError("PDF 文件不存在")
    header = path.read_bytes()[:5]
    if not header.startswith(b"%PDF-"):
        raise ValueError(
            "文件不是有效的 PDF；若由 PPT 转换，请正确导出 PDF 或直接上传 .pptx"
        )

    blocks: list[ParsedBlock] = []
    try:
        with pdfplumber.open(path) as pdf:
            for idx, page in enumerate(pdf.pages, start=1):
                text = (page.extract_text() or "").strip()
                if text and _is_meaningful_text(text):
                    blocks.append(ParsedBlock(text=text, page=idx))
    except Exception as exc:
        raise ValueError(
            f"PDF 解析失败：{exc}。"
            "若由 PPT 转换而来，请重新导出 PDF 或直接上传 .pptx 文件。"
        ) from exc

    if not blocks:
        raise ValueError(PDF_PARSE_EMPTY_MSG)
    return blocks


def _collect_shape_text(shape, parts: list[str]) -> None:
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
    except ImportError:
        MSO_SHAPE_TYPE = None

    if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        for child in shape.shapes:
            _collect_shape_text(child, parts)
        return

    if getattr(shape, "has_text_frame", False):
        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if text:
                parts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))


def _parse_pptx(path: Path) -> list[ParsedBlock]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise ValueError("服务器未安装 python-pptx，无法解析 PPTX") from exc

    if not path.is_file():
        raise ValueError("PPTX 文件不存在")
    if path.read_bytes()[:2] != b"PK":
        raise ValueError("文件不是有效的 PPTX（Office Open XML）格式")

    prs = Presentation(str(path))
    blocks: list[ParsedBlock] = []
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            _collect_shape_text(shape, parts)

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(notes)

        slide_text = "\n".join(dict.fromkeys(parts)).strip()
        if slide_text and _is_meaningful_text(slide_text):
            blocks.append(ParsedBlock(text=slide_text, page=idx))

    if not blocks:
        raise ValueError(PPTX_PARSE_EMPTY_MSG)
    return blocks
